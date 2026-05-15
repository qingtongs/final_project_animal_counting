from __future__ import annotations

import json
from collections import Counter
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSET_DIR = ROOT / "ppt_assets"
OUT = ROOT / "project_presentation_draft.pptx"

W, H = 13.333, 7.5
NAVY = RGBColor(23, 37, 84)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(22, 101, 52)
AMBER = RGBColor(180, 83, 9)
SLATE = RGBColor(51, 65, 85)
LIGHT = RGBColor(248, 250, 252)
MINT = RGBColor(220, 252, 231)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(185, 28, 28)


def font(size=28, bold=False, color=SLATE):
    return {"size": Pt(size), "bold": bold, "color": color}


def set_run(run, spec):
    run.font.name = "Aptos"
    run.font.size = spec["size"]
    run.font.bold = spec["bold"]
    run.font.color.rgb = spec["color"]


def add_text(slide, text, x, y, w, h, spec=None, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    set_run(p.runs[0], spec or font())
    return shape


def add_bullets(slide, items, x, y, w, h, size=18, color=SLATE):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        set_run(p.runs[0], font(size, False, color))
    return shape


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.55, 0.35, 8.6, 0.45, font(24, True, NAVY))
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.86, 9.5, 0.35, font(11, False, SLATE))


def add_card(slide, x, y, w, h, title, value=None, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(226, 232, 240)
    add_text(slide, title, x + 0.18, y + 0.15, w - 0.35, 0.28, font(10, True, SLATE))
    if value is not None:
        add_text(slide, value, x + 0.18, y + 0.52, w - 0.35, h - 0.62, font(26, True, accent))


def cover_image(path: Path, out: Path, size=(1600, 900), darken=True):
    img = Image.open(path).convert("RGB")
    img_ratio = img.width / img.height
    target_ratio = size[0] / size[1]
    if img_ratio > target_ratio:
        new_w = int(img.height * target_ratio)
        left = (img.width - new_w) // 2
        img = img.crop((left, 0, left + new_w, img.height))
    else:
        new_h = int(img.width / target_ratio)
        top = (img.height - new_h) // 2
        img = img.crop((0, top, img.width, top + new_h))
    img = img.resize(size, Image.Resampling.LANCZOS)
    if darken:
        overlay = Image.new("RGB", size, (0, 0, 0))
        img = Image.blend(img, overlay, 0.35)
    img.save(out)


def get_font(size=28, bold=False):
    candidates = [
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/msyh.ttc",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size)
    return ImageFont.load_default()


def draw_bar_chart(data: dict[str, int], out: Path, title: str, color=(37, 99, 235)):
    width, height = 1300, 620
    img = Image.new("RGB", (width, height), "white")
    d = ImageDraw.Draw(img)
    title_font = get_font(34, True)
    label_font = get_font(19)
    small_font = get_font(17)
    d.text((50, 30), title, fill=(23, 37, 84), font=title_font)
    items = sorted(data.items(), key=lambda kv: (-kv[1], kv[0]))
    max_v = max(data.values()) if data else 1
    x0, y0, chart_w, chart_h = 80, 115, 1120, 390
    bar_gap = 8
    bar_w = max(18, (chart_w - bar_gap * (len(items) - 1)) // len(items))
    d.line((x0, y0 + chart_h, x0 + chart_w, y0 + chart_h), fill=(148, 163, 184), width=2)
    for i, (name, value) in enumerate(items):
        x = x0 + i * (bar_w + bar_gap)
        bh = int(chart_h * value / max_v)
        y = y0 + chart_h - bh
        d.rounded_rectangle((x, y, x + bar_w, y0 + chart_h), radius=7, fill=color)
        d.text((x + bar_w / 2, y - 24), str(value), fill=(15, 23, 42), font=small_font, anchor="mm")
        label = name[:8]
        d.text((x + bar_w / 2, y0 + chart_h + 26), label, fill=(51, 65, 85), font=label_font, anchor="mm")
    img.save(out)


def draw_scores(scores: dict[str, float], out: Path):
    draw_bar_chart({k.replace(".png", ""): round(v) for k, v in scores.items()}, out, "Zero-shot YOLO-World baseline score by validation image", (22, 101, 52))


def draw_sample_grid(out: Path):
    files = ["val_1.png", "val_4.png", "val_7.png", "val_9.png"]
    labels = json.loads((ROOT / "val_set" / "ground_truth.json").read_text(encoding="utf-8"))
    cell_w, cell_h = 520, 390
    img = Image.new("RGB", (cell_w * 2, cell_h * 2), (248, 250, 252))
    d = ImageDraw.Draw(img)
    font_b = get_font(22, True)
    font_s = get_font(17)
    for idx, name in enumerate(files):
        src = Image.open(ROOT / "val_set" / name).convert("RGB")
        src.thumbnail((cell_w - 26, cell_h - 88), Image.Resampling.LANCZOS)
        x = (idx % 2) * cell_w + 13
        y = (idx // 2) * cell_h + 13
        img.paste(src, (x, y))
        text_y = y + src.height + 8
        d.text((x, text_y), name, fill=(23, 37, 84), font=font_b)
        label = ", ".join(f"{k}:{v}" for k, v in labels[name].items())
        d.text((x, text_y + 28), label, fill=(51, 65, 85), font=font_s)
    img.save(out)


def score_sample(truth, pred):
    n = len(truth)
    score = 0.0
    for label, count in truth.items():
        if label in pred:
            score += 50 / n
            if pred[label] == count:
                score += 50 / n
    score -= 5 * len(set(pred) - set(truth))
    return max(score, 0.0)


def build_assets():
    ASSET_DIR.mkdir(exist_ok=True)
    cover_image(ROOT / "val_set" / "val_7.png", ASSET_DIR / "cover.png")
    draw_sample_grid(ASSET_DIR / "sample_grid.png")
    truth = json.loads((ROOT / "val_set" / "ground_truth.json").read_text(encoding="utf-8"))
    freq = Counter()
    totals = {}
    for name, labels in truth.items():
        totals[name.replace(".png", "")] = sum(labels.values())
        freq.update(labels)
    draw_bar_chart(dict(freq), ASSET_DIR / "category_frequency.png", "Validation-set animal category frequency", (37, 99, 235))
    draw_bar_chart(totals, ASSET_DIR / "animals_per_image.png", "Total animals per validation image", (180, 83, 9))
    pred_path = ROOT / "detect_smoke_predictions.json"
    if pred_path.exists():
        preds = json.loads(pred_path.read_text(encoding="utf-8"))
        scores = {name: score_sample(labels, preds.get(name, {})) for name, labels in truth.items()}
    else:
        scores = {}
    draw_scores(scores or {k: 0 for k in truth}, ASSET_DIR / "baseline_scores.png")
    return truth, scores


def add_footer(slide, idx):
    add_text(slide, f"{idx:02d}", 12.25, 7.02, 0.5, 0.22, font(9, False, RGBColor(100, 116, 139)), PP_ALIGN.RIGHT)


def build_deck():
    truth, scores = build_assets()
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # 1 Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(ASSET_DIR / "cover.png"), 0, 0, width=Inches(W), height=Inches(H))
    add_text(slide, "Animal Detection\nand Counting", 0.75, 1.0, 6.3, 1.5, font(42, True, WHITE))
    add_text(slide, "YOLO-based multimodal label dictionary extraction", 0.82, 2.82, 7.0, 0.45, font(19, False, MINT))
    add_text(slide, "Deep Learning Final Project | 2025-2026 Spring", 0.82, 6.45, 6.0, 0.35, font(13, False, WHITE))

    # 2 Task
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Project task and evaluation target", "Convert each image into a clean machine-readable animal-count dictionary.")
    add_card(slide, 0.65, 1.45, 2.4, 1.2, "Input", "one image", BLUE)
    add_card(slide, 3.35, 1.45, 2.4, 1.2, "Output", "{label: count}", GREEN)
    add_card(slide, 6.05, 1.45, 2.4, 1.2, "Constraint", "animal only", AMBER)
    add_card(slide, 8.75, 1.45, 2.4, 1.2, "Score", "category + count", RED)
    add_bullets(slide, [
        "Only categories observed in the image are written to JSON.",
        "Counts must be positive integers; zero-count categories are omitted.",
        "Extra wrong categories receive a penalty, so precision matters.",
    ], 0.85, 3.15, 8.8, 1.7, 21)

    # 3 Dataset examples
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Validation images are dense multi-species scenes", "Synthetic images contain repeated instances, similar species, and background distractions.")
    slide.shapes.add_picture(str(ASSET_DIR / "sample_grid.png"), Inches(0.68), Inches(1.25), width=Inches(7.1))
    add_bullets(slide, [
        "Examples cover same-class counting and multi-class recognition.",
        "Scenes include partial overlap and visually similar animal shapes.",
        "Validation labels are used only for local format and scorer checks.",
    ], 8.15, 1.55, 4.35, 3.2, 18)

    # 4 Category frequency
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Validation distribution guides threshold tuning", "Frequent classes expose counting risks; rare classes expose recall risks.")
    slide.shapes.add_picture(str(ASSET_DIR / "category_frequency.png"), Inches(0.8), Inches(1.25), width=Inches(11.7))

    # 5 Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "YOLO counting pipeline", "The code separates detection, filtering, counting, and JSON formatting.")
    steps = ["Load image", "YOLO-World\ninference", "Confidence\nfiltering", "NMS duplicate\nhandling", "Animal whitelist", "Count + JSON"]
    xs = [0.65, 2.65, 4.65, 6.65, 8.65, 10.65]
    for i, (x, step) in enumerate(zip(xs, steps)):
        add_card(slide, x, 2.2, 1.45, 1.15, f"Step {i+1}", step, BLUE if i < 2 else GREEN)
        if i < len(xs) - 1:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 1.48), Inches(2.78), Inches(x + 1.9), Inches(2.78))
            line.line.color.rgb = RGBColor(148, 163, 184)
            line.line.width = Pt(2)
    add_bullets(slide, [
        "Model: yolov8s-worldv2.pt through Ultralytics YOLOWorld.",
        "Classes are set to the 20 required animal labels before inference.",
        "Final JSON contains no natural-language explanations.",
    ], 1.0, 4.35, 10.8, 1.35, 19)

    # 6 Training strategy
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Training and fine-tuning strategy", "The current code runs zero-shot; the next improvement is supervised fine-tuning on category-balanced synthetic data.")
    add_bullets(slide, [
        "Data sources: public animal datasets, web-collected images, and synthetic multi-animal compositions.",
        "Annotation: bounding boxes in YOLO format with the 20-class whitelist.",
        "Augmentation: scale jitter, horizontal flip, color jitter, mosaic, random occlusion.",
        "Validation: tune confidence and IoU thresholds against dictionary-level score, not only mAP.",
    ], 0.8, 1.35, 6.1, 4.4, 19)
    add_card(slide, 7.35, 1.42, 2.2, 1.1, "Backbone", "YOLO-World", BLUE)
    add_card(slide, 9.85, 1.42, 2.2, 1.1, "Classes", "20 animals", GREEN)
    add_card(slide, 7.35, 2.95, 2.2, 1.1, "Metric", "mAP + JSON", AMBER)
    add_card(slide, 9.85, 2.95, 2.2, 1.1, "Goal", "fewer extras", RED)

    # 7 Results
    mean_score = sum(scores.values()) / len(scores) if scores else 0.0
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Current experiment results", "A real smoke test shows the zero-shot model runs end-to-end, but class confusion remains.")
    add_card(slide, 0.75, 1.35, 2.6, 1.25, "Format-check oracle", "100.00", GREEN)
    add_card(slide, 3.65, 1.35, 2.6, 1.25, "Zero-shot baseline", f"{mean_score:.2f}", AMBER)
    add_card(slide, 6.55, 1.35, 2.6, 1.25, "Images processed", "10 / 10", BLUE)
    add_card(slide, 9.45, 1.35, 2.6, 1.25, "Next target", "fine-tune", RED)
    add_bullets(slide, [
        "100.00 is the validation-format check using the released ground-truth JSON.",
        f"{mean_score:.2f} is the actual zero-shot YOLO-World baseline score from detect_smoke_predictions.json.",
        "Fine-tuning should reduce false sheep/zebra/dog predictions and improve recall for lion, tiger, fox, monkey, and pig.",
    ], 0.9, 3.25, 10.9, 1.9, 18)

    # 8 Baseline chart
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Baseline score varies by scene complexity", "The lowest-scoring images combine multiple species and visually similar bodies.")
    slide.shapes.add_picture(str(ASSET_DIR / "baseline_scores.png"), Inches(0.8), Inches(1.18), width=Inches(11.7))

    # 9 Example output comparison
    preds = json.loads((ROOT / "detect_smoke_predictions.json").read_text(encoding="utf-8"))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Example: prediction dictionary comparison", "The failure pattern is not JSON format; it is animal class discrimination.")
    slide.shapes.add_picture(str(ROOT / "val_set" / "val_1.png"), Inches(0.75), Inches(1.25), width=Inches(4.6), height=Inches(4.6))
    add_text(slide, "Ground truth", 5.75, 1.35, 2.4, 0.35, font(18, True, GREEN))
    add_bullets(slide, [json.dumps(truth["val_1.png"], ensure_ascii=False)], 5.75, 1.85, 3.2, 1.0, 18)
    add_text(slide, "Zero-shot output", 5.75, 3.05, 2.8, 0.35, font(18, True, AMBER))
    add_bullets(slide, [json.dumps(preds["val_1.png"], ensure_ascii=False)], 5.75, 3.55, 3.2, 1.0, 18)
    add_bullets(slide, [
        "Correct: cat and goose counts.",
        "Missed: two lions and one fox.",
        "False extra: sheep detections from similar body regions.",
    ], 9.2, 1.65, 3.25, 3.0, 17)

    # 10 Error analysis
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Error analysis and remedies", "The scoring rule rewards both category recall and exact counts, while penalizing extra categories.")
    add_card(slide, 0.75, 1.35, 3.0, 1.35, "Main error", "class confusion", RED)
    add_card(slide, 4.05, 1.35, 3.0, 1.35, "Counting risk", "overlap / occlusion", AMBER)
    add_card(slide, 7.35, 1.35, 3.0, 1.35, "Fix", "fine-tune + tune", GREEN)
    add_bullets(slide, [
        "Duck/goose/chicken require examples with similar pose and white feathers.",
        "Fox/wolf/dog and lion/tiger/cat require harder negatives.",
        "Threshold selection should optimize the project score, because extra categories cost 5 points each.",
        "Class-balanced synthetic scenes can target rare labels such as monkey, goat, giraffe, and deer.",
    ], 0.95, 3.3, 10.7, 2.2, 19)

    # 11 Submission workflow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Final evaluation workflow", "The evaluation day script is intentionally simple and reproducible.")
    add_bullets(slide, [
        "Install dependencies: python -m pip install -r requirements.txt",
        "Run evaluation: python evaluate.py --image-dir eval_set --output predictions.json",
        "Check JSON validity before submission.",
        "Package report, code ZIP, and presentation with the required naming rule.",
    ], 0.85, 1.35, 7.0, 3.4, 21)
    slide.shapes.add_picture(str(ASSET_DIR / "animals_per_image.png"), Inches(7.5), Inches(1.35), width=Inches(5.1))

    # 12 Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Conclusion", "The system is ready for batch JSON generation; accuracy improves next through targeted fine-tuning.")
    add_bullets(slide, [
        "Completed: animal whitelist, robust JSON writer, validation scorer, and batch evaluation script.",
        "Verified: validation format score is 100.00, and zero-shot inference runs end-to-end.",
        "Next: collect/annotate class-balanced data, fine-tune YOLO, and tune thresholds on dictionary-level score.",
    ], 1.0, 1.6, 10.8, 3.0, 23)

    for i, slide in enumerate(prs.slides, 1):
        add_footer(slide, i)

    prs.save(OUT)
    print(OUT)
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    build_deck()
