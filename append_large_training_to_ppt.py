from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
RUN = ROOT / "runs" / "animal_openimages3000_yolov8n_e20_gpu"
PPT = ROOT / "project_presentation_draft.pptx"


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=RGBColor(30, 41, 59)):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = shape.text_frame.paragraphs[0]
    p.text = text
    r = p.runs[0]
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return shape


def card(slide, x, y, w, h, title, value, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(226, 232, 240)
    add_text(slide, title, x + 0.12, y + 0.12, w - 0.24, 0.25, 10, True)
    add_text(slide, value, x + 0.12, y + 0.48, w - 0.24, 0.42, 20, True, color)


def main() -> None:
    s = json.loads((ROOT / "training_summary_large.json").read_text(encoding="utf-8"))
    prs = Presentation(PPT)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    add_text(slide, "Larger, higher-quality training set", 0.55, 0.35, 8.5, 0.45, 25, True, RGBColor(23, 37, 84))
    add_text(slide, "Open Images V7 subset expanded from 600 to 3000 training images, with real bbox annotations and 5858 animal boxes.", 0.58, 0.85, 11.5, 0.3, 12)
    card(slide, 0.7, 1.35, 2.0, 0.95, "Train images", str(s["train_images"]), RGBColor(37, 99, 235))
    card(slide, 2.95, 1.35, 2.0, 0.95, "Train boxes", str(s["train_boxes"]), RGBColor(22, 101, 52))
    card(slide, 5.2, 1.35, 2.0, 0.95, "Val images", str(s["val_images"]), RGBColor(180, 83, 9))
    card(slide, 7.45, 1.35, 2.0, 0.95, "Epochs", str(s["epochs"]), RGBColor(185, 28, 28))
    card(slide, 9.7, 1.35, 2.0, 0.95, "Classes", "19 + wolf empty", RGBColor(23, 37, 84))
    slide.shapes.add_picture(str(RUN / "labels.jpg"), Inches(0.75), Inches(2.62), width=Inches(5.4))
    slide.shapes.add_picture(str(RUN / "train_batch0.jpg"), Inches(6.75), Inches(2.62), width=Inches(5.4))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    add_text(slide, "Large training run: real metrics", 0.55, 0.35, 7.5, 0.45, 25, True, RGBColor(23, 37, 84))
    card(slide, 0.7, 1.25, 2.0, 0.95, "Precision", f"{s['precision']:.3f}", RGBColor(37, 99, 235))
    card(slide, 2.95, 1.25, 2.0, 0.95, "Recall", f"{s['recall']:.3f}", RGBColor(22, 101, 52))
    card(slide, 5.2, 1.25, 2.0, 0.95, "mAP50", f"{s['map50']:.3f}", RGBColor(180, 83, 9))
    card(slide, 7.45, 1.25, 2.0, 0.95, "mAP50-95", f"{s['map5095']:.3f}", RGBColor(185, 28, 28))
    card(slide, 9.7, 1.25, 2.0, 0.95, "Project val", "35.67", RGBColor(23, 37, 84))
    slide.shapes.add_picture(str(RUN / "results.png"), Inches(0.7), Inches(2.55), width=Inches(5.8))
    slide.shapes.add_picture(str(RUN / "confusion_matrix_normalized.png"), Inches(6.9), Inches(2.55), width=Inches(5.25))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    add_text(slide, "Training scale improved results", 0.55, 0.35, 7.5, 0.45, 25, True, RGBColor(23, 37, 84))
    add_text(slide, "Measured comparison across the small Open Images run, the larger Open Images run, and YOLO-World on the synthetic validation set.", 0.58, 0.85, 11.5, 0.3, 12)
    card(slide, 0.8, 1.55, 3.2, 1.2, "Small Open Images", "mAP50 0.297\nscore 19.58", RGBColor(185, 28, 28))
    card(slide, 4.25, 1.55, 3.2, 1.2, "Large Open Images", "mAP50 0.504\nscore 35.67", RGBColor(22, 101, 52))
    card(slide, 7.7, 1.55, 3.2, 1.2, "YOLO-World", "score 37.58", RGBColor(37, 99, 235))
    add_text(slide, "Interpretation", 0.85, 3.35, 2.2, 0.35, 19, True, RGBColor(23, 37, 84))
    add_text(slide, "The larger dataset substantially improves general animal detection and narrows the gap on the synthetic validation images. The remaining gap is mainly domain style: the final project images contain dense, generated multi-animal compositions that differ from natural Open Images photos.", 0.85, 3.85, 10.8, 1.05, 20)
    slide.shapes.add_picture(str(RUN / "val_batch0_pred.jpg"), Inches(0.85), Inches(5.2), width=Inches(5.25))
    slide.shapes.add_picture(str(RUN / "val_batch0_labels.jpg"), Inches(6.55), Inches(5.2), width=Inches(5.25))

    prs.save(PPT)
    print(PPT)


if __name__ == "__main__":
    main()
