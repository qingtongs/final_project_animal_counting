from __future__ import annotations

import csv
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
RUN = ROOT / "runs" / "animal_openimages_yolov8n_e10_gpu"
PPT = ROOT / "project_presentation_draft.pptx"


def add_text(slide, text, x, y, w, h, size=20, bold=False, color=RGBColor(30, 41, 59)):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = shape.text_frame.paragraphs[0]
    p.text = text
    r = p.runs[0]
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return shape


def card(slide, x, y, w, h, title, value, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(226, 232, 240)
    add_text(slide, title, x + 0.12, y + 0.12, w - 0.24, 0.25, 10, True)
    add_text(slide, value, x + 0.12, y + 0.48, w - 0.24, 0.45, 22, True, color)


def latest_metrics():
    with (RUN / "results.csv").open("r", encoding="utf-8") as f:
        rows = list(csv.DictReader(f))
    row = rows[-1]
    return {
        "precision": float(row["metrics/precision(B)"]),
        "recall": float(row["metrics/recall(B)"]),
        "map50": float(row["metrics/mAP50(B)"]),
        "map5095": float(row["metrics/mAP50-95(B)"]),
        "train_box_loss": float(row["train/box_loss"]),
        "train_cls_loss": float(row["train/cls_loss"]),
        "train_dfl_loss": float(row["train/dfl_loss"]),
    }


def main():
    metrics = latest_metrics()
    summary = {
        **metrics,
        "dataset": "Open Images V7 subset exported to YOLO format",
        "train_images": 600,
        "val_images": 150,
        "epochs": 10,
        "model": "YOLOv8n pretrained -> 20-class animal head",
        "weights": str(RUN / "weights" / "best.pt"),
        "project_val_best_conf": 0.05,
        "project_val_score_after_training": 19.58,
        "project_val_score_yoloworld_zero_shot": 37.58,
        "note": "Fine-tuning improved Open Images bbox metrics but did not improve the synthetic validation-set dictionary score because of domain shift.",
    }
    (ROOT / "training_summary.json").write_text(json.dumps(summary, indent=2), encoding="utf-8")

    prs = Presentation(PPT)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    add_text(slide, "Real supervised training result", 0.55, 0.35, 7.5, 0.45, 25, True, RGBColor(23, 37, 84))
    add_text(slide, "Dataset: 600 train / 150 validation Open Images V7 samples, YOLOv8n, 10 epochs, RTX 3070 Ti Laptop GPU.", 0.58, 0.85, 11.5, 0.3, 12)
    card(slide, 0.65, 1.35, 2.1, 0.95, "Precision", f"{metrics['precision']:.3f}", RGBColor(37, 99, 235))
    card(slide, 2.95, 1.35, 2.1, 0.95, "Recall", f"{metrics['recall']:.3f}", RGBColor(22, 101, 52))
    card(slide, 5.25, 1.35, 2.1, 0.95, "mAP50", f"{metrics['map50']:.3f}", RGBColor(180, 83, 9))
    card(slide, 7.55, 1.35, 2.1, 0.95, "mAP50-95", f"{metrics['map5095']:.3f}", RGBColor(185, 28, 28))
    card(slide, 9.85, 1.35, 2.1, 0.95, "Best weights", "best.pt", RGBColor(23, 37, 84))
    slide.shapes.add_picture(str(RUN / "results.png"), Inches(0.75), Inches(2.65), width=Inches(5.8))
    slide.shapes.add_picture(str(RUN / "confusion_matrix_normalized.png"), Inches(7.0), Inches(2.65), width=Inches(5.2))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    add_text(slide, "Project validation after training", 0.55, 0.35, 7.5, 0.45, 25, True, RGBColor(23, 37, 84))
    add_text(slide, "The trained detector is real, but the teacher-provided synthetic validation set has a domain gap from Open Images.", 0.58, 0.85, 11.4, 0.3, 12)
    card(slide, 0.75, 1.45, 2.6, 1.1, "YOLO-World baseline", "37.58", RGBColor(22, 101, 52))
    card(slide, 3.65, 1.45, 2.6, 1.1, "Trained YOLO best", "19.58", RGBColor(185, 28, 28))
    card(slide, 6.55, 1.45, 2.6, 1.1, "Best conf", "0.05", RGBColor(180, 83, 9))
    card(slide, 9.45, 1.45, 2.6, 1.1, "Format score", "100.00", RGBColor(37, 99, 235))
    add_text(slide, "Interpretation", 0.8, 3.05, 2.2, 0.3, 18, True, RGBColor(23, 37, 84))
    add_text(
        slide,
        "Open Images fine-tuning learns real bbox labels, but the validation scenes are synthetic collages with multiple animals and unusual co-occurrence. For the final hidden synthetic set, YOLO-World remains the stronger default unless a synthetic-style annotated dataset is added.",
        0.8,
        3.55,
        11.1,
        1.2,
        19,
    )
    slide.shapes.add_picture(str(RUN / "val_batch0_pred.jpg"), Inches(0.8), Inches(5.05), width=Inches(5.3))
    slide.shapes.add_picture(str(RUN / "val_batch0_labels.jpg"), Inches(6.55), Inches(5.05), width=Inches(5.3))

    prs.save(PPT)
    print(PPT)
    print(ROOT / "training_summary.json")


if __name__ == "__main__":
    main()
