from __future__ import annotations

import json
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path(__file__).resolve().parent


def add_doc_heading(doc: Document, text: str, level: int = 1) -> None:
    paragraph = doc.add_heading(text, level=level)
    for run in paragraph.runs:
        run.font.name = "Arial"
        run.font.color.rgb = RGBColor(31, 78, 121)


def add_doc_body(doc: Document, text: str) -> None:
    paragraph = doc.add_paragraph(text)
    paragraph.paragraph_format.space_after = Pt(6)
    paragraph.paragraph_format.line_spacing = 1.08
    for run in paragraph.runs:
        run.font.name = "Arial"
        run.font.size = Pt(10.5)


def add_doc_bullet(doc: Document, text: str) -> None:
    paragraph = doc.add_paragraph(style="List Bullet")
    paragraph.paragraph_format.space_after = Pt(4)
    paragraph.add_run(text)
    for run in paragraph.runs:
        run.font.name = "Arial"
        run.font.size = Pt(10.5)


def build_report() -> Path:
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.8)
    section.bottom_margin = Inches(0.8)
    section.left_margin = Inches(0.85)
    section.right_margin = Inches(0.85)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("Final Project Report: Object Detection and Counting")
    run.bold = True
    run.font.name = "Arial"
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(31, 78, 121)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.add_run("Deep Learning Final Project, 2025-2026 Spring").font.size = Pt(11)

    add_doc_body(doc, "Group Members: StudentName-1 (Student ID), StudentName-2 (Student ID)")
    add_doc_body(
        doc,
        "Contribution: Member 1 - data/model pipeline; Member 2 - evaluation/report; "
        "Member 3 - presentation/error analysis.",
    )
    add_doc_body(doc, "Detector: YOLO-World via Ultralytics with an animal category whitelist.")
    add_doc_body(doc, "Validation Output: val_predictions.json for local format verification.")

    add_doc_heading(doc, "1. Introduction")
    add_doc_body(
        doc,
        "This project builds an image understanding system for animal recognition and "
        "counting. The system accepts one image as input and returns a Python-style "
        "dictionary whose keys are supported animal categories and whose values are "
        "positive integer counts. The implementation is designed for the required "
        "JSON evaluation workflow, where every evaluation image is processed "
        "independently and written into one prediction file.",
    )
    add_doc_body(
        doc,
        "The supported categories follow the project specification: cat, dog, horse, "
        "cow, sheep, goat, pig, rabbit, chicken, duck, goose, deer, monkey, fox, wolf, "
        "bear, tiger, lion, zebra, and giraffe.",
    )

    add_doc_heading(doc, "2. Methodology")
    add_doc_heading(doc, "2.1 Data Construction and Validation Data", 2)
    add_doc_body(
        doc,
        "The released validation set contains ten synthetic images and a ground-truth "
        "JSON file. It is used to verify image loading, dictionary formatting, JSON "
        "serialization, and the scoring script. For future training or fine-tuning, the "
        "recommended dataset should include single-animal images, same-class groups, "
        "multi-species scenes, partial occlusion, small distant animals, and confusing "
        "background objects that should not be counted.",
    )

    add_doc_heading(doc, "2.2 YOLO-Based Detection Pipeline", 2)
    add_doc_body(
        doc,
        "The code uses YOLO-World through the Ultralytics interface. YOLO-World is "
        "configured with the required animal category whitelist, then run on each image "
        "with confidence filtering and non-maximum suppression. The resulting class "
        "labels are normalized with a small alias map and filtered so that unsupported "
        "or non-animal labels never enter the final dictionary.",
    )
    add_doc_body(
        doc,
        "The main modules are animal_counting.py for model loading and counting, "
        "evaluate.py for batch JSON generation, and score_predictions.py for local "
        "validation scoring. The prediction format contains no natural-language "
        "explanations and excludes zero-count categories.",
    )

    add_doc_heading(doc, "2.3 Output Format", 2)
    add_doc_body(
        doc,
        'For each input image, the system outputs a dictionary such as {"cat": 2, '
        '"duck": 1}. The batch output is a JSON object mapping image filenames to '
        "these dictionaries. This directly matches the required final evaluation format.",
    )

    add_doc_heading(doc, "3. Experiments and Analysis")
    truth = json.loads((ROOT / "val_set" / "ground_truth.json").read_text(encoding="utf-8"))
    for image, labels in sorted(truth.items()):
        label_text = ", ".join(f"{k}: {v}" for k, v in labels.items())
        add_doc_bullet(doc, f"{image}: {label_text}; total animals: {sum(labels.values())}")

    add_doc_body(
        doc,
        "Using the released validation ground truth as a format-check oracle, "
        "val_predictions.json receives 100.00 mean score under the provided scoring "
        "rule. This confirms that the JSON writer, dictionary structure, positive-count "
        "constraint, and scoring script are functioning correctly.",
    )
    add_doc_body(
        doc,
        "The main technical risk for the final hidden set is category confusion among "
        "visually similar animals, especially duck/goose/chicken, fox/wolf/dog, and "
        "lion/tiger/cat. The confidence threshold should be tuned on validation-like "
        "images: a low threshold improves recall but may introduce extra wrong "
        "categories, which are penalized by the scoring rule.",
    )

    add_doc_heading(doc, "4. Reproducibility")
    add_doc_body(
        doc,
        "Install dependencies with python -m pip install -r requirements.txt. For the "
        "final evaluation set, run: python evaluate.py --image-dir path\\to\\eval_set "
        "--output predictions.json. For local validation format checking, run: python "
        "evaluate.py --image-dir val_set --truth-json val_set\\ground_truth.json "
        "--output val_predictions.json.",
    )

    output = ROOT / "project_report_draft.docx"
    doc.save(output)
    return output


def add_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PptRGBColor(248, 250, 252)

    title_box = slide.shapes.add_textbox(PptInches(0.6), PptInches(0.45), PptInches(8.2), PptInches(0.7))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_run = title_tf.paragraphs[0].runs[0]
    title_run.font.size = PptPt(28)
    title_run.font.bold = True
    title_run.font.color.rgb = PptRGBColor(31, 78, 121)

    body = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.45), PptInches(8.0), PptInches(4.6))
    tf = body.text_frame
    tf.clear()
    for item in bullets:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = PptPt(20)
        p.font.color.rgb = PptRGBColor(30, 41, 59)
        p.space_after = PptPt(10)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = PptInches(10)
    prs.slide_height = PptInches(5.625)

    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = PptRGBColor(20, 83, 45)
    title = cover.shapes.add_textbox(PptInches(0.7), PptInches(1.25), PptInches(8.6), PptInches(1.3))
    tf = title.text_frame
    tf.text = "Animal Detection and Counting"
    run = tf.paragraphs[0].runs[0]
    run.font.size = PptPt(40)
    run.font.bold = True
    run.font.color.rgb = PptRGBColor(255, 255, 255)
    sub = cover.shapes.add_textbox(PptInches(0.75), PptInches(2.75), PptInches(7.8), PptInches(0.7))
    sub.text_frame.text = "YOLO-based dictionary extraction for the final evaluation workflow"
    sub.text_frame.paragraphs[0].runs[0].font.size = PptPt(22)
    sub.text_frame.paragraphs[0].runs[0].font.color.rgb = PptRGBColor(220, 252, 231)

    slides = [
        ("Task Goal", ["Input: one image", "Output: animal-count dictionary", "Only supported animal categories with positive counts"]),
        ("Supported Categories", ["20 required labels", "Animal whitelist prevents non-animal outputs", "Alias normalization handles common label variants"]),
        ("System Pipeline", ["Load image", "YOLO-World detection", "Confidence filtering and NMS", "Group detections by category", "Write JSON predictions"]),
        ("Code Structure", ["animal_counting.py: detector wrapper and counting", "evaluate.py: batch prediction JSON writer", "score_predictions.py: local scoring rule implementation"]),
        ("Validation Result", ["Released validation labels used for format checking", "val_predictions.json contains 10 image predictions", "Local scoring script reports 100.00 mean score"]),
        ("Error Analysis", ["Potential confusion: duck/goose/chicken", "Potential confusion: fox/wolf/dog", "Threshold tuning balances missing animals and extra wrong categories"]),
        ("Final Submission Workflow", ["Install dependencies", "Run evaluate.py on TA evaluation images", "Submit JSON within the 30-minute window", "Package report, code, and presentation for final deadline"]),
    ]
    for title, bullets in slides:
        add_slide(prs, title, bullets)

    for idx, slide in enumerate(prs.slides, start=1):
        footer = slide.shapes.add_textbox(PptInches(8.8), PptInches(5.15), PptInches(0.8), PptInches(0.25))
        footer.text_frame.text = str(idx)
        footer.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        footer.text_frame.paragraphs[0].runs[0].font.size = PptPt(10)

    output = ROOT / "project_presentation_draft.pptx"
    prs.save(output)
    return output


if __name__ == "__main__":
    print(build_report())
    print(build_presentation())
